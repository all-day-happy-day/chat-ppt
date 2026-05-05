from pathlib import Path
from typing import cast

from pptx import Presentation as PPT
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide, SlideLayout
from ulid import ULID

from app.bible.domain.repository import BibleRepository
from app.bible.domain.valueobject import BiblePhrase
from app.config import config
from app.powerpoint.domain.entity import Layout
from app.powerpoint.domain.repository import LayoutRepository
from app.powerpoint.domain.service import TemplateReadService
from app.project.domain.entity import BiblePart, LyricsPart, PlainPart, ValuePart
from app.project.domain.service import PresentationService
from app.project.domain.types import Parts
from app.shared.bible.domain.enum import AvailableBibleVersions


class PPTXPresentationService(PresentationService):
    LYRICS_LINES_PER_SLIDE: int = 2

    def __init__(
        self,
        bible_repository: BibleRepository,
        layout_repository: LayoutRepository,
        template_read_service: TemplateReadService,
    ) -> None:
        self.bible_repository: BibleRepository = bible_repository
        self.layout_repository: LayoutRepository = layout_repository
        self.template_read_service: TemplateReadService = template_read_service

    def _write_plain(self, ppt: Presentation, layouts_dict: dict[str, SlideLayout], part: PlainPart) -> None:
        if part.layout_id is None:
            raise ValueError("Layout ID is required")
        layout_entity: Layout = self.layout_repository.get_by_id(id=part.layout_id)
        ppt.slides.add_slide(layouts_dict[layout_entity.name])

    def _write_value(
        self,
        ppt: Presentation,
        layouts_dict: dict[str, SlideLayout],
        layout_id: ULID,
        placeholder_content_dict: dict[int, str],
    ) -> None:
        layout_entity: Layout = self.layout_repository.get_by_id(id=layout_id)
        slide: Slide = ppt.slides.add_slide(layouts_dict[layout_entity.name])

        # ppt_placeholders: SlidePlaceholders = slide.placeholders
        placeholders: list[BaseShape] = [
            ppt_shape for ppt_shape in slide.slide_layout.shapes if ppt_shape.is_placeholder
        ]
        ppt_placeholders: list[BaseShape] = [shape for shape in slide.placeholders]

        for ph, ppt_ph in zip(placeholders, ppt_placeholders):
            if ph.shape_id in placeholder_content_dict:
                cast(AutoShape, ppt_ph).text = placeholder_content_dict[ph.shape_id]

    def _write_lyrics(self, ppt: Presentation, layouts_dict: dict[str, SlideLayout], part: LyricsPart) -> None:
        if part.lyrics_layout_id is None:
            raise ValueError("Lyrics layout ID is required")
        lyrics_layout_entity: Layout = self.layout_repository.get_by_id(id=part.lyrics_layout_id)

        title_layout_entity: Layout | None
        if part.title_layout_id is None:
            title_layout_entity = None
        else:
            title_layout_entity = self.layout_repository.get_by_id(id=part.title_layout_id)

        for content in part.contents.contents:
            if content.include_title_slide and title_layout_entity:
                # write title
                title_placeholder_shape_id: int | None = part.contents.title_placeholder_shape_id
                if title_placeholder_shape_id is None:
                    raise ValueError("Title placeholder shape ID is required")
                self._write_value(
                    ppt=ppt,
                    layouts_dict=layouts_dict,
                    layout_id=title_layout_entity.id,
                    placeholder_content_dict={title_placeholder_shape_id: content.title},
                )

            # write lyrics
            for song_part_idx in content.lyrics_part_sequence:
                lyrics_text: str
                if song_part_idx == -1:
                    lyrics_text = ""
                else:
                    lyrics_text = content.lyrics[song_part_idx].lyrics
                lyrics_lines: list[str] = lyrics_text.split("\n")
                for lyrics_line_idx in range(0, len(lyrics_lines), self.LYRICS_LINES_PER_SLIDE):
                    lyrics_lines_slice: list[str] = lyrics_lines[
                        lyrics_line_idx : lyrics_line_idx + self.LYRICS_LINES_PER_SLIDE
                    ]
                    lyrics_line_text: str = "\n".join(lyrics_lines_slice)
                    # TODO: Line overflow check and fallback
                    self._write_value(
                        ppt=ppt,
                        layouts_dict=layouts_dict,
                        layout_id=lyrics_layout_entity.id,
                        placeholder_content_dict={part.contents.lyrics_placeholder_shape_id: lyrics_line_text},
                    )

    def _write_bible(self, ppt: Presentation, layouts_dict: dict[str, SlideLayout], part: BiblePart) -> None:
        if part.phrase_layout_id is None:
            raise ValueError("Phrase layout ID is required")
        phrase_layout_entity: Layout = self.layout_repository.get_by_id(id=part.phrase_layout_id)

        title_layout_entity: Layout | None
        if part.title_layout_id is None:
            title_layout_entity = None
        else:
            title_layout_entity = self.layout_repository.get_by_id(id=part.title_layout_id)

        for phrase_range in part.contents.contents:
            if phrase_range.type == "phrase":
                assert phrase_range.start is not None
                version: AvailableBibleVersions = phrase_range.start.version
                book: str = phrase_range.start.book
                chapter: int = phrase_range.start.chapter
                verse: int = phrase_range.start.verse

                end_verse: int
                if phrase_range.end is not None:
                    end_verse = phrase_range.end.verse
                else:
                    end_verse = verse

                for verse in range(verse, end_verse + 1):
                    phrase: BiblePhrase = self.bible_repository.get(
                        version=version, book=book, chapter=chapter, verse=verse
                    )
                    # TODO: Slide overflow check and fallback
                    placeholder_content_dict: dict[int, str] = {part.contents.phrase_placeholder_id: phrase.phrase}
                    if part.contents.phrase_range_placeholder_id is not None:
                        placeholder_content_dict[part.contents.phrase_range_placeholder_id] = (
                            f"{book} {chapter}:{verse}"
                        )
                    self._write_value(
                        ppt=ppt,
                        layouts_dict=layouts_dict,
                        layout_id=phrase_layout_entity.id,
                        placeholder_content_dict=placeholder_content_dict,
                    )
            elif phrase_range.type == "title" and title_layout_entity:
                # write title
                self._write_value(
                    ppt=ppt,
                    layouts_dict=layouts_dict,
                    layout_id=title_layout_entity.id,
                    placeholder_content_dict=part.contents.title_placeholder_values or {},
                )

    def generate_and_save(
        self,
        template_file_path: str | Path,
        parts: list[Parts],
        save_ppt_filename: str,
        project_id: ULID,
        user_id: ULID,
    ) -> Path:
        ppt: Presentation = PPT(pptx=str(template_file_path))
        layouts: list[SlideLayout] = self.template_read_service._get_template_slide_layouts(ppt=ppt)
        layouts_dict: dict[str, SlideLayout] = {layout.name: layout for layout in layouts}

        parts.sort(key=lambda part: part.order)
        for part in parts:
            if isinstance(part, PlainPart):
                self._write_plain(ppt=ppt, layouts_dict=layouts_dict, part=part)
            elif isinstance(part, ValuePart):
                layout_id: ULID | None = part.layout_id
                if layout_id is None:
                    raise ValueError("Layout ID is required")
                placeholder_content_dict: dict[int, str] = {
                    content.placeholder_shape_id: content.value or "" for content in part.contents.contents
                }

                self._write_value(
                    ppt=ppt,
                    layouts_dict=layouts_dict,
                    layout_id=layout_id,
                    placeholder_content_dict=placeholder_content_dict,
                )
            elif isinstance(part, LyricsPart):
                self._write_lyrics(ppt=ppt, layouts_dict=layouts_dict, part=part)
            elif isinstance(part, BiblePart):
                self._write_bible(ppt=ppt, layouts_dict=layouts_dict, part=part)
            else:
                raise ValueError(f"Invalid part type: {type(part)}")

        if Path(save_ppt_filename).suffix == ".pptx":
            save_ppt_filename = str(Path(save_ppt_filename).stem)
        save_path: Path = Path(config.ppt_save_temp_directory) / f"p{project_id}_u{user_id} / {save_ppt_filename}.pptx"
        save_path.parent.mkdir(parents=True, exist_ok=True)
        ppt.save(str(save_path))

        return save_path
