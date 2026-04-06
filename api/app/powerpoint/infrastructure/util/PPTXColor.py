from typing import Self

from bs4 import BeautifulSoup, Tag
from pptx.dml.color import ColorFormat
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.slide import SlideLayout
from pydantic import BaseModel
from pydantic_extra_types import Color

from app.powerpoint.domain.valueobject import ColorConfig


class _ThemeColor(BaseModel):
    accent1: str
    accent2: str
    accent3: str
    accent4: str
    accent5: str
    accent6: str
    bg1: str
    bg2: str
    dk1: str
    dk2: str
    hlink: str
    followed_hyperlink: str
    lt1: str
    lt2: str
    tx1: str
    tx2: str

    @classmethod
    def from_slide_layout(cls, slide_layout: SlideLayout) -> Self:
        color_map_tag: Tag | None = BeautifulSoup(
            markup=slide_layout.slide_master.part.blob.decode("utf-8"), features="xml"
        ).find("p:clrMap")
        assert color_map_tag is not None
        color_map: dict[str, str] = {k: str(v) for k, v in color_map_tag.attrs.items()}

        theme_part: str = slide_layout.slide_master.part.part_related_by(reltype=RT.THEME).blob.decode("utf-8")
        theme_soup: BeautifulSoup = BeautifulSoup(markup=theme_part, features="xml")

        xml_colors: dict[str, str] = {
            "lt1": theme_soup.find("a:clrScheme").find("a:lt1").find("a:sysClr")["lastClr"],  # type: ignore
            "lt2": theme_soup.find("a:clrScheme").find("a:lt2").find("a:srgbClr")["val"],  # type: ignore
            "dk1": theme_soup.find("a:clrScheme").find("a:dk1").find("a:sysClr")["lastClr"],  # type: ignore
            "dk2": theme_soup.find("a:clrScheme").find("a:dk2").find("a:srgbClr")["val"],  # type: ignore
            "accent1": theme_soup.find("a:clrScheme").find("a:accent1").find("a:srgbClr")["val"],  # type: ignore
            "accent2": theme_soup.find("a:clrScheme").find("a:accent2").find("a:srgbClr")["val"],  # type: ignore
            "accent3": theme_soup.find("a:clrScheme").find("a:accent3").find("a:srgbClr")["val"],  # type: ignore
            "accent4": theme_soup.find("a:clrScheme").find("a:accent4").find("a:srgbClr")["val"],  # type: ignore
            "accent5": theme_soup.find("a:clrScheme").find("a:accent5").find("a:srgbClr")["val"],  # type: ignore
            "accent6": theme_soup.find("a:clrScheme").find("a:accent6").find("a:srgbClr")["val"],  # type: ignore
            "hlink": theme_soup.find("a:clrScheme").find("a:hlink").find("a:srgbClr")["val"],  # type: ignore
            "folHlink": theme_soup.find("a:clrScheme").find("a:folHlink").find("a:srgbClr")["val"],  # type: ignore
        }

        return cls(
            accent1=xml_colors[color_map["accent1"]],
            accent2=xml_colors[color_map["accent2"]],
            accent3=xml_colors[color_map["accent3"]],
            accent4=xml_colors[color_map["accent4"]],
            accent5=xml_colors[color_map["accent5"]],
            accent6=xml_colors[color_map["accent6"]],
            bg1=xml_colors[color_map["bg1"]],
            bg2=xml_colors[color_map["bg2"]],
            dk1=xml_colors["dk1"],
            dk2=xml_colors["dk2"],
            hlink=xml_colors[color_map["hlink"]],
            followed_hyperlink=xml_colors[color_map["folHlink"]],
            lt1=xml_colors["lt1"],
            lt2=xml_colors["lt2"],
            tx1=xml_colors[color_map["tx1"]],
            tx2=xml_colors[color_map["tx2"]],
        )


def _solid_color_resolver(color: ColorFormat, theme_color: _ThemeColor | None = None) -> tuple[str, float]:
    color_code: str
    if color.type == MSO_COLOR_TYPE.RGB:
        color_code = str(color.rgb)
    elif color.type == MSO_COLOR_TYPE.SCHEME:
        if theme_color is None:
            raise ValueError("Theme color is not defined")
        color_code = getattr(theme_color, color.theme_color.xml_value)
        if color_code is None:
            raise ValueError(f"Theme color {color.theme_color.xml_value} is not defined")
    else:
        raise ValueError(f"Invalid color type: {color.type}")

    alpha: float
    alpha_element: Tag | None = BeautifulSoup(markup=color._color._xClr.xml, features="xml").find("a:alpha")
    if alpha_element is None:
        alpha = 1.0
    else:
        alpha = int(str(alpha_element["val"])) / 100000.0

    return color_code, alpha


def resolve_layout_background_color(slide_layout: SlideLayout) -> ColorConfig:
    theme_color: _ThemeColor = _ThemeColor.from_slide_layout(slide_layout=slide_layout)

    # Get Layout Background Color
    if slide_layout.background.fill.type == MSO_FILL_TYPE.SOLID:
        layout_background_color_fill: ColorFormat = slide_layout.background.fill.fore_color
        layout_background_color_code, _ = _solid_color_resolver(
            color=layout_background_color_fill, theme_color=theme_color
        )
        return ColorConfig(color_type="solid", color=Color(layout_background_color_code), alpha=1.0)
    elif slide_layout.background.fill.type == MSO_FILL_TYPE.BACKGROUND:
        pass
    else:
        return ColorConfig(color_type="none", color=None, alpha=None)

    # Get Theme Background Color
    if slide_layout.background.fill.type == MSO_FILL_TYPE.SOLID:
        theme_background_color_fill: ColorFormat = slide_layout.background.fill.fore_color
        theme_background_color_code, _ = _solid_color_resolver(
            color=theme_background_color_fill, theme_color=theme_color
        )
        return ColorConfig(color_type="solid", color=Color(theme_background_color_code), alpha=1.0)
    else:
        return ColorConfig(color_type="none", color=None, alpha=None)


def resolve_shape_fill_color(shape: BaseShape, theme_color: _ThemeColor) -> ColorConfig:
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        assert isinstance(shape, AutoShape)
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            color_code, color_alpha = _solid_color_resolver(color=shape.fill.fore_color, theme_color=theme_color)
            return ColorConfig(color_type="solid", color=Color(color_code), alpha=color_alpha)
        else:
            return ColorConfig(color_type="none", color=None, alpha=None)
    else:
        return ColorConfig(color_type="none", color=None, alpha=None)
