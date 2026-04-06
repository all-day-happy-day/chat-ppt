from datetime import datetime, timezone
from pathlib import Path
from typing import cast

from pptx import Presentation as PPT
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.slide import SlideLayout
from pptx.util import Length
from ulid import ULID

from app.powerpoint.domain.entity import Layout, Shape, Template, TemplateFile
from app.powerpoint.domain.enum import ShapeType
from app.powerpoint.domain.exception import InvalidFileReadRequest, SlideSizeNotDefined
from app.powerpoint.domain.service.TemplateReadService import TemplateReadService
from app.powerpoint.domain.valueobject import ColorConfig, Position, Size
from app.powerpoint.infrastructure.util.PPTXColor import (
    _ThemeColor,
    resolve_layout_background_color,
    resolve_shape_fill_color,
)


class PPTXTemplateReadService(TemplateReadService):
    theme_color: _ThemeColor | None = None

    def _get_template_slide_layouts(self, ppt: Presentation) -> list[SlideLayout]:
        """Gather all slide layouts from the template.
        There could be multiple slide masters in the template, but we gather all slide layouts from all slide masters.
        """
        slide_layouts: list[SlideLayout] = []
        for slide_master in ppt.slide_masters:
            for slide_layout in slide_master.slide_layouts:
                slide_layouts.append(slide_layout)
        return slide_layouts

    def _get_slide_size(self, ppt: Presentation) -> Size:
        sw: Length | None = ppt.slide_width
        sh: Length | None = ppt.slide_height
        if sw is None or sh is None:
            raise SlideSizeNotDefined("Slide size is not defined")
        sw_emu: int = cast(Length, sw).emu
        sh_emu: int = cast(Length, sh).emu
        return Size(width=sw_emu, height=sh_emu)

    def _get_shape_type(self, shape: BaseShape) -> ShapeType:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return ShapeType.IMAGE
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return ShapeType.SHAPE
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return ShapeType.TEXT
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return ShapeType.TABLE
        else:
            return ShapeType.OTHER

    def _get_shape_element(
        self, shape: GroupShape | BaseShape, layout_id: ULID, group: GroupShape | None = None
    ) -> list[Shape]:
        shape_list: list[Shape] = []

        if isinstance(shape, GroupShape):
            for child_shape in shape.shapes:
                shape_list.extend(self._get_shape_element(shape=child_shape, layout_id=layout_id, group=shape))
        elif isinstance(shape, BaseShape):
            shape_size: Size = Size(width=shape.width.emu, height=shape.height.emu)
            shape_position: Position = Position(x=shape.left.emu, y=shape.top.emu, rotation=shape.rotation)

            if group is not None:
                sx: float = group.width.emu / group.element.chExt.cx.emu if group.element.chExt.cx.emu != 0 else 1.0
                sy: float = group.height.emu / group.element.chExt.cy.emu if group.element.chExt.cy.emu != 0 else 1.0
                shape_position = Position(
                    x=group.left.emu + round((shape.left.emu - group.element.chOff.x.emu) * sx),
                    y=group.top.emu + round((shape.top.emu - group.element.chOff.y.emu) * sy),
                    rotation=shape.rotation,
                )
                shape_size = Size(width=round(shape_size.width * sx), height=round(shape_size.height * sy))

            if self.theme_color is None:
                raise ValueError("Theme color is not defined")
            fill_color: ColorConfig = resolve_shape_fill_color(shape=shape, theme_color=self.theme_color)

            shape_list.append(
                Shape(
                    id=ULID(),
                    layout_id=layout_id,
                    shape_id=shape.shape_id,
                    size=shape_size,
                    position=shape_position,
                    name=shape.name,
                    text=cast(AutoShape, shape).text if shape.has_text_frame else None,  # TODO: check neeeded
                    placeholder=shape.is_placeholder,
                    type=self._get_shape_type(shape=shape),
                    fill_color=fill_color,
                )
            )
        else:
            raise ValueError(f"Invalid shape type: {type(shape)}")
        return shape_list

    def _get_layout(self, slide_layout: SlideLayout, template_id: ULID) -> Layout:
        layout_id: ULID = ULID()
        shape_list: list[Shape] = []
        self.theme_color = _ThemeColor.from_slide_layout(slide_layout=slide_layout)
        for shape in slide_layout.shapes:
            shape_list.extend(self._get_shape_element(shape=shape, layout_id=layout_id))
        self.theme_color = None
        return Layout(
            id=layout_id,
            template_id=template_id,
            name=slide_layout.name,
            background_color=resolve_layout_background_color(slide_layout=slide_layout),
            shapes=shape_list,
        )

    def _get_template(
        self,
        ppt: Presentation,
        user_id: ULID,
        template_name: str | None = None,
        template: Template | None = None,
    ) -> Template:
        template_id = template.id if template is not None else ULID()
        slide_layouts: list[SlideLayout] = self._get_template_slide_layouts(ppt=ppt)
        layouts: list[Layout] = [
            self._get_layout(slide_layout=slide_layout, template_id=template_id) for slide_layout in slide_layouts
        ]
        now: datetime = datetime.now(tz=timezone.utc)
        if template is not None:
            return template.model_copy(
                update={
                    "name": template_name if template_name is not None else template.name,
                    "slide_size": self._get_slide_size(ppt=ppt),
                    # "updated_at": now,
                    "layouts": layouts,
                }
            )
        else:
            if template_name is None:
                raise InvalidFileReadRequest("Template name is required")
            return Template(
                id=template_id,
                user_id=user_id,
                name=template_name,
                slide_size=self._get_slide_size(ppt=ppt),
                created_at=now,
                updated_at=now,
                layouts=layouts,
            )

    def read(
        self, user_id: ULID, ppt_path: str | Path, template_name: str | None = None, template: Template | None = None
    ) -> tuple[TemplateFile, Template]:
        ppt: Presentation = PPT(pptx=str(ppt_path))
        template = self._get_template(ppt=ppt, template_name=template_name, user_id=user_id, template=template)

        template_file: TemplateFile = TemplateFile(
            template_id=template.id,
            user_id=user_id,
            path=str(ppt_path),
            size=Path(ppt_path).stat().st_size,
        )
        return template_file, template
